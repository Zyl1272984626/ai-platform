from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


OUT = Path(r"C:\Users\Administrator\Desktop\demo_editable_16x9.pptx")


BLUE = RGBColor(31, 81, 171)
TEAL = RGBColor(22, 124, 138)
GREEN = RGBColor(59, 126, 55)
ORANGE = RGBColor(233, 103, 30)
DARK = RGBColor(25, 34, 45)
GRAY = RGBColor(125, 128, 132)
LIGHT_BLUE = RGBColor(242, 247, 255)
LIGHT_TEAL = RGBColor(239, 251, 252)
LIGHT_GREEN = RGBColor(242, 250, 241)
LIGHT_ORANGE = RGBColor(255, 247, 240)
LIGHT_PURPLE = RGBColor(248, 245, 255)
WHITE = RGBColor(255, 255, 255)


def rgb(hex_value):
    hex_value = hex_value.lstrip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def set_line(shape, color=BLUE, width=1.2, dash=None):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    if dash:
        shape.line.dash_style = dash


def add_textbox(slide, x, y, w, h, text, font_size=12, color=DARK, bold=False,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font="Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def format_run(run, font_size, color=DARK, bold=False, font="Microsoft YaHei"):
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_formula_text(slide, x, y, w, h, lines, font_size=11, color=DARK, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.text = line
        for run in p.runs:
            format_run(run, font_size, color, bold, "Cambria Math")
    return box


def add_box(slide, x, y, w, h, text="", line=BLUE, fill=WHITE, radius=True, font_size=11,
            text_color=DARK, bold=False, dash=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    set_line(shp, line, 1.1, dash)
    if text:
        tf = shp.text_frame
        tf.clear()
        tf.margin_left = Inches(0.04)
        tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = text
        for run in p.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(font_size)
            run.font.color.rgb = text_color
            run.font.bold = bold
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=BLUE, width=2.0, dash=None, begin=False, end=True):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    if dash:
        ln.line.dash_style = dash
    line_xml = ln.line._get_or_add_ln()
    for child in list(line_xml):
        if child.tag.endswith("}headEnd") or child.tag.endswith("}tailEnd"):
            line_xml.remove(child)
    if begin:
        head = OxmlElement("a:headEnd")
        head.set("type", "triangle")
        head.set("w", "med")
        head.set("len", "med")
        line_xml.append(head)
    if end:
        tail = OxmlElement("a:tailEnd")
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")
        line_xml.append(tail)
    return ln


def add_down_arrow(slide, x, y1, y2, color=BLUE, width=1.8):
    return add_arrow(slide, x, y1, x, y2, color=color, width=width)


def add_block_arrow(slide, x, y, w, h, color=BLUE, direction="right"):
    shape_map = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
        "up_down": MSO_SHAPE.UP_DOWN_ARROW,
    }
    shp = slide.shapes.add_shape(shape_map[direction], Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = color
    return shp


def add_icon_cloud(slide, x, y, color=BLUE):
    add_box(slide, x + 0.02, y + 0.22, 0.34, 0.12, "", color, WHITE, radius=True)
    for dx, dy, s in [(0.00, 0.12, 0.18), (0.13, 0.06, 0.24), (0.32, 0.13, 0.18)]:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + dx), Inches(y + dy), Inches(s), Inches(s))
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        set_line(c, color, 1.2)
    add_arrow(slide, x + 0.13, y + 0.42, x + 0.13, y + 0.60, color, 1.1, end=False)
    add_arrow(slide, x + 0.28, y + 0.42, x + 0.28, y + 0.60, color, 1.1, end=False)
    add_arrow(slide, x + 0.03, y + 0.60, x + 0.38, y + 0.60, color, 1.1, end=False)


def add_icon_people(slide, x, y, color=ORANGE):
    for dx in [0.0, 0.32]:
        head = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + dx), Inches(y), Inches(0.16), Inches(0.16))
        head.fill.solid()
        head.fill.fore_color.rgb = WHITE
        set_line(head, color, 1.1)
        body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + dx - 0.02), Inches(y + 0.18), Inches(0.20), Inches(0.26))
        body.fill.solid()
        body.fill.fore_color.rgb = WHITE
        set_line(body, color, 1.1)
    bubble = slide.shapes.add_shape(MSO_SHAPE.CLOUD_CALLOUT, Inches(x + 0.17), Inches(y - 0.08), Inches(0.26), Inches(0.20))
    bubble.fill.solid()
    bubble.fill.fore_color.rgb = WHITE
    set_line(bubble, color, 1.0)


def add_icon_clipboard(slide, x, y, color=GREEN):
    board = add_box(slide, x, y, 0.36, 0.50, "", color, WHITE, radius=True)
    add_box(slide, x + 0.11, y - 0.04, 0.14, 0.10, "", color, WHITE, radius=True)
    for i in range(3):
        add_arrow(slide, x + 0.10, y + 0.15 + i * 0.11, x + 0.26, y + 0.15 + i * 0.11, color, 0.9, end=False)
    check = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(x + 0.30), Inches(y + 0.36), Inches(0.22), Inches(0.22))
    check.fill.background()
    set_line(check, color, 1.0)


def slide_one(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 4.82, 0.12, 3.70, 0.35, "图3-1  课题三总体架构图", 18, DARK, True)
    add_textbox(slide, 4.18, 0.47, 4.95, 0.25, "面向身心健康的智能心理干预系统：感知-决策-干预闭环架构", 10, RGBColor(72, 78, 86))

    add_box(slide, 0.25, 0.88, 0.95, 4.52, "课题一：\n可信感知\n与监测", GRAY, rgb("f7f7f7"), True, 12, DARK, True)
    add_box(slide, 3.25, 0.88, 0.95, 4.52, "课题二：\n智能理疗\n决策", GRAY, rgb("f7f7f7"), True, 12, DARK, True)

    left_items = [("情绪压力+睡眠质量+\n代谢风险三维状态估计", 1.65),
                  ("置信度+预测区间", 2.85),
                  ("趋势预测+异常预警", 4.05)]
    for text, y in left_items:
        add_box(slide, 1.62, y, 1.45, 0.55, text, BLUE, WHITE, True, 8.5, DARK, True)
        add_block_arrow(slide, 1.22, y + 0.19, 0.36, 0.18, BLUE, "right")

    mid_items = [("中医证型辨识+\n置信度分布", 1.68),
                 ("结构化电刺激处方\n(穴位组合+波形+\n频率+强度+时长)", 2.72),
                 ("自适应调控策略", 4.05)]
    for text, y in mid_items:
        add_box(slide, 4.55, y, 1.35, 0.57, text, ORANGE, WHITE, True, 8.2, DARK, True)
        add_block_arrow(slide, 4.22, y + 0.20, 0.32, 0.18, ORANGE, "right")

    outer = add_box(slide, 6.10, 0.90, 6.30, 4.50, "", TEAL, WHITE, True)
    add_textbox(slide, 7.55, 1.02, 3.25, 0.25, "课题三： 智能心理干预系统与创新应用", 11, TEAL, True)

    add_box(slide, 6.20, 1.35, 6.10, 1.15, "", TEAL, LIGHT_TEAL, True)
    add_textbox(slide, 6.30, 1.48, 1.00, 0.25, "云端AI服务", 10, TEAL, True)
    cloud_boxes = [("多模态感知\n推理引擎", 7.34), ("证型适配\n决策引擎", 8.58), ("大模型智能体\n对话引擎", 9.82), ("在线策略\n学习引擎", 11.18)]
    for text, x in cloud_boxes:
        add_box(slide, x, 1.52, 0.96, 0.60, text, TEAL, WHITE, True, 7.7, DARK, True)
    for x1, x2 in [(8.30, 8.58), (9.54, 9.82), (10.78, 11.18)]:
        add_arrow(slide, x1, 1.82, x2, 1.82, GRAY, 1.0, begin=True, end=True)
    for x in [7.82, 9.05, 10.30, 11.62]:
        add_down_arrow(slide, x, 2.40, 2.13, TEAL, 1.2)
    add_arrow(slide, 7.82, 2.40, 11.62, 2.40, TEAL, 1.2, dash=MSO_LINE_DASH_STYLE.DASH, end=False)
    add_textbox(slide, 8.70, 2.54, 1.30, 0.25, "WiFi/4G", 10, TEAL, True)
    add_block_arrow(slide, 8.54, 2.56, 0.16, 0.28, TEAL, "up_down")

    add_box(slide, 6.20, 2.92, 6.10, 0.80, "", BLUE, LIGHT_BLUE, True)
    add_textbox(slide, 6.30, 3.16, 0.95, 0.25, "移动端App", 10, BLUE, True)
    app_boxes = [("身心状态\n仪表盘", 7.36), ("干预控制\n面板", 8.34), ("智能体\n对话界面", 9.25), ("情绪日记", 10.25), ("中医\n知识库", 11.15)]
    for text, x in app_boxes:
        add_box(slide, x, 3.05, 0.72, 0.48, text, BLUE, WHITE, True, 7.5, DARK, True)
    add_textbox(slide, 8.72, 3.80, 1.05, 0.25, "BLE 5.0", 10, TEAL, True)
    add_block_arrow(slide, 8.54, 3.70, 0.16, 0.32, TEAL, "up_down")

    add_box(slide, 6.20, 4.14, 6.10, 1.10, "", GREEN, LIGHT_GREEN, True)
    add_textbox(slide, 6.30, 4.33, 0.80, 0.25, "智能腕带", 10, GREEN, True)
    add_box(slide, 7.15, 4.36, 0.50, 0.62, "♡\n⌁", GRAY, WHITE, True, 12, DARK, True)
    band_items = [("PPG/EDA/ST/IMU四模态传感器", 8.00, 4.25, 2.10),
                  ("PC6/HT7双通道电刺激模组", 10.22, 4.25, 1.80),
                  ("端侧推理引擎", 8.00, 4.72, 2.10),
                  ("端侧安全校验", 10.22, 4.72, 1.80)]
    for text, x, y, w in band_items:
        add_box(slide, x, y, w, 0.34, text, GREEN, WHITE, True, 7.8, DARK, True)

    add_arrow(slide, 12.42, 3.37, 12.12, 3.37, BLUE, 1.5)
    add_textbox(slide, 12.55, 4.68, 0.60, 0.50, "真实用户\n数据反馈", 10, BLUE, True)
    add_arrow(slide, 12.72, 4.60, 12.72, 3.44, BLUE, 1.2, dash=MSO_LINE_DASH_STYLE.DASH, end=True)
    add_arrow(slide, 12.72, 5.24, 0.64, 5.24, BLUE, 1.2, dash=MSO_LINE_DASH_STYLE.DASH, end=True)
    add_arrow(slide, 0.64, 5.38, 0.64, 5.24, BLUE, 1.2)
    add_arrow(slide, 3.72, 5.36, 3.72, 5.24, BLUE, 1.2)

    add_box(slide, 0.85, 6.00, 3.15, 0.80, "", BLUE, WHITE, True, dash=MSO_LINE_DASH_STYLE.DASH)
    add_icon_cloud(slide, 1.07, 6.10, BLUE)
    add_textbox(slide, 1.82, 6.15, 1.90, 0.45, "研究内容一：\n系统设计与端边云协同", 10, BLUE, True)
    add_block_arrow(slide, 4.30, 6.29, 0.50, 0.22, BLUE, "right")

    add_box(slide, 5.08, 6.00, 3.20, 0.80, "", ORANGE, WHITE, True, dash=MSO_LINE_DASH_STYLE.DASH)
    add_icon_people(slide, 5.28, 6.16, ORANGE)
    add_textbox(slide, 6.00, 6.14, 1.95, 0.48, "研究内容二：\n多场景干预方案与协同交互", 10, ORANGE, True)
    add_block_arrow(slide, 8.50, 6.29, 0.50, 0.22, BLUE, "right")

    add_box(slide, 9.12, 6.00, 3.20, 0.80, "", GREEN, WHITE, True, dash=MSO_LINE_DASH_STYLE.DASH)
    add_icon_clipboard(slide, 9.40, 6.15, GREEN)
    add_textbox(slide, 10.05, 6.14, 1.95, 0.48, "研究内容三：\n真实用户验证与迭代优化", 10, GREEN, True)
    add_textbox(slide, 4.82, 7.10, 3.70, 0.25, "设计集成 → 场景交互 → 验证优化", 13, BLUE, True)


def slide_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 1.68, 0.10, 10.70, 1.18, "", BLUE, rgb("fbfdff"), True)
    add_textbox(slide, 1.85, 0.42, 1.15, 0.43, "1. 输入层\n(多模态数据)", 11, BLUE, True)
    inputs = [("SQE 信号质量评估\n+ 预处理", 3.05), ("PPG/HRV\nPPG: 光电容积脉搏波", 5.12), ("EDA\nEDA: 皮电活动", 7.04), ("ST\nST: 皮肤温度", 8.70), ("IMU\nIMU: 惯性测量单元", 10.35)]
    icons = ["▮▮▮", "♥", "●", "♨", "↗"]
    for idx, (text, x) in enumerate(inputs):
        add_box(slide, x, 0.18, 1.78, 0.74, "", BLUE, WHITE, True)
        add_textbox(slide, x + 0.14, 0.37, 0.25, 0.20, icons[idx], 13, BLUE, True)
        add_textbox(slide, x + 0.42, 0.30, 1.18, 0.38, text, 8.0, DARK, True)
        add_down_arrow(slide, x + 0.89, 0.92, 1.18, DARK, 1.0)
    add_arrow(slide, 3.95, 1.18, 11.23, 1.18, DARK, 1.0, end=True)

    add_box(slide, 0.15, 0.78, 1.12, 5.86, "", BLUE, WHITE, True, 11, BLUE, True, dash=MSO_LINE_DASH_STYLE.DASH)
    add_textbox(slide, 0.32, 1.02, 0.78, 0.25, "训练策略", 11, BLUE, True)
    steps = [("①\nWESAD /\nDEAP /\nSHHS /\nMESA\n公开数据\n预训练", 1.45),
             ("②\n本项目\n真实佩戴\n数据微调", 4.00),
             ("③\nEMA /\nPSG\n参考校准", 5.52)]
    for text, y in steps:
        add_textbox(slide, 0.32, y, 0.78, 1.25, text, 10, DARK, True)
    add_down_arrow(slide, 0.72, 3.32, 3.75, DARK, 1.8)
    add_down_arrow(slide, 0.72, 5.03, 5.38, DARK, 1.8)

    add_box(slide, 1.68, 1.48, 10.70, 2.86, "", BLUE, rgb("fbfdff"), True)
    add_textbox(slide, 1.85, 1.62, 1.45, 0.25, "2. 解耦表征层", 12, BLUE, True)
    add_textbox(slide, 5.45, 1.62, 3.10, 0.25, "语义解耦表征学习（β-VAE主干）", 12, BLUE, True)
    add_box(slide, 1.84, 1.92, 1.70, 1.69, "模态专属编码器\n\nPPG/HRV 编码器\nEDA 编码器\nST 编码器\nIMU 编码器\n……", BLUE, WHITE, True, 9, DARK, True)
    add_arrow(slide, 3.62, 2.62, 3.90, 2.62, DARK, 2.0)
    add_box(slide, 3.95, 1.95, 1.62, 1.30, "共享表征 zshared\n\n自主神经负荷、\n昼夜节律等\n跨模态信息", BLUE, LIGHT_BLUE, True, 8.5, DARK, True)
    add_box(slide, 6.58, 1.95, 1.55, 1.30, "特有表征 zspecific\n\n模态独有\n生理含义", GRAY, WHITE, True, 8.5, DARK, True)
    add_box(slide, 9.14, 1.95, 1.40, 1.30, "噪声表征 znoise\n\n运动伪影、\n佩戴干扰", GRAY, rgb("f7f7f7"), True, 8.5, DARK, True)
    for y, label in [(2.34, "互信息约束\nLMI"), (2.85, "正交约束\nLorth")]:
        add_arrow(slide, 5.60, y, 6.48, y, DARK, 1.0, dash=MSO_LINE_DASH_STYLE.DASH, begin=True, end=True)
        add_textbox(slide, 5.73, y - 0.32, 0.62, 0.25, label, 8.5, DARK, False)
    for y, label in [(2.34, "互信息约束\nLMI"), (2.85, "独立性约束")]:
        add_arrow(slide, 8.17, y, 9.05, y, DARK, 1.0, dash=MSO_LINE_DASH_STYLE.DASH, begin=True, end=True)
        add_textbox(slide, 8.30, y - 0.32, 0.62, 0.25, label, 8.5, DARK, False)
    add_box(slide, 3.95, 3.43, 6.55, 0.78, "", BLUE, WHITE, True, dash=MSO_LINE_DASH_STYLE.DASH)
    add_textbox(slide, 5.38, 3.47, 1.70, 0.20, "缺失鲁棒性训练", 10, BLUE, True)
    add_textbox(slide, 4.60, 3.68, 2.20, 0.22, "随机模态遮挡  →  重构共享语义", 8.5, DARK, False)
    add_formula_text(slide, 4.72, 3.88, 1.95, 0.20, ["L_miss = ∑||s_m − ŝ'_m||²"], 10, DARK)
    add_arrow(slide, 7.78, 3.45, 7.78, 4.20, BLUE, 0.9, dash=MSO_LINE_DASH_STYLE.DASH, end=False)
    add_textbox(slide, 8.25, 3.72, 1.75, 0.42, "关键模态缺失时\n传递低置信标记", 9, DARK, False)
    add_arrow(slide, 10.86, 1.62, 10.86, 4.08, BLUE, 0.9, dash=MSO_LINE_DASH_STYLE.DASH, end=False)
    add_box(slide, 11.05, 1.65, 1.45, 2.36, "", BLUE, WHITE, True)
    add_formula_text(
        slide, 11.25, 1.88, 1.05, 1.88,
        [
            "L_total = L_rec",
            "+ λ₁L_KL",
            "+ λ₂L_MI",
            "+ λ₃L_Orth",
            "+ λ₄L_Miss",
            "+ λ₅L_Phy",
            "+ λ₆L_N",
            "+ λ₇L_T",
        ],
        11,
        DARK,
    )
    add_down_arrow(slide, 6.88, 4.34, 4.62, DARK, 1.8)

    add_box(slide, 1.68, 4.50, 10.70, 1.65, "", BLUE, rgb("fbfdff"), True)
    add_textbox(slide, 1.86, 4.62, 1.90, 0.25, "3. 先验约束融合层", 12, BLUE, True)
    add_textbox(slide, 5.20, 4.62, 3.62, 0.25, "候选时滞依赖发现与先验约束融合", 12, BLUE, True)
    add_box(slide, 2.18, 4.88, 2.35, 1.10, "时间尺度分解\n秒级（应激反应）\n分钟级（情绪波动）\n小时级（睡眠节律）", BLUE, WHITE, True, 9, DARK, True)
    add_arrow(slide, 4.55, 5.42, 4.90, 5.42, DARK, 1.8)
    add_box(slide, 4.95, 4.88, 1.95, 1.10, "因果发现\nNeural Granger +\nPCMCI\n→ 候选时滞结构", BLUE, WHITE, True, 9, DARK, False)
    add_arrow(slide, 6.90, 5.42, 7.28, 5.42, DARK, 1.8)
    add_box(slide, 7.35, 4.88, 3.00, 1.10, "", BLUE, WHITE, True)
    add_textbox(slide, 8.22, 4.98, 1.25, 0.20, "先验约束融合", 9, BLUE, True)
    add_formula_text(slide, 7.68, 5.24, 2.35, 0.22, ["αᵢⱼ = softmax(Qᵢ · Kⱼ + β · Δₜₛ(i,j))"], 9.2, DARK)
    add_textbox(slide, 8.02, 5.64, 1.70, 0.20, "生理先验筛除不合理连接", 8.5, DARK, False)
    add_arrow(slide, 10.84, 5.42, 10.35, 5.42, DARK, 1.0, dash=MSO_LINE_DASH_STYLE.DASH)
    add_box(slide, 10.86, 4.88, 1.70, 1.10, "生理先验\n交感-副交感平衡、\n体温-睡眠节律、\n心率-运动-能耗", BLUE, WHITE, True, 8.5, DARK, True)
    add_down_arrow(slide, 6.88, 6.15, 6.38, DARK, 1.8)

    add_box(slide, 1.68, 6.30, 10.70, 0.78, "", BLUE, LIGHT_PURPLE, True)
    add_textbox(slide, 1.86, 6.42, 1.05, 0.25, "4. 输出层", 12, BLUE, True)
    add_textbox(slide, 5.35, 6.39, 3.05, 0.25, "多模态融合表征输出", 12, rgb("4938a3"), True)
    add_box(slide, 3.12, 6.70, 7.86, 0.34, "向研究内容二输入  →  ht（融合表征向量）", rgb("7a66d6"), WHITE, True, 9.5, DARK, True)
    add_textbox(slide, 5.00, 7.18, 3.45, 0.28, "图2 研究内容一技术路线图", 16, DARK, True)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    slide_one(prs)
    slide_two(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
